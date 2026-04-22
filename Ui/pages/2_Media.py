"""
🎙️ Media Analyzer — Upload Audio/Video & Analyze

Cyberpunk-themed interface for uploading meeting recordings (audio/video),
transcribing them with local Whisper, extracting role-aware intelligence
via Gemma 4, and exporting summaries to PDF or PPTX.
"""

from __future__ import annotations

import io
import os
import sys
import tempfile
from pathlib import Path
from datetime import datetime

import pandas as pd
import streamlit as st
import streamlit.components.v1 as components

# ── Ensure project root is importable ────────────────────────────────
PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from core.audio import AudioProcessor, ALL_EXTENSIONS  # noqa: E402
from core.engine import MeetingAnalyzer  # noqa: E402
from core.storage import save_analysis  # noqa: E402

# ── Add Ui directory to path for vault_sidebar import ────────────────
UI_DIR = Path(__file__).resolve().parent.parent
if str(UI_DIR) not in sys.path:
    sys.path.insert(0, str(UI_DIR))

from vault_sidebar import render_vault_sidebar  # noqa: E402


# ── Cached resources (survive reruns) ────────────────────────────────
@st.cache_resource
def get_analyzer() -> MeetingAnalyzer:
    return MeetingAnalyzer()


@st.cache_resource
def get_audio_processor() -> AudioProcessor:
    return AudioProcessor()


# ── Export helpers ────────────────────────────────────────────────────
def _build_pdf(result, role: str, filename: str) -> bytes:
    """Generate a styled PDF report from MeetingOutput using fpdf2."""
    try:
        from fpdf import FPDF
    except ImportError:
        raise ImportError("fpdf2 is required for PDF export. Run: pip install fpdf2")

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    # Title
    pdf.set_font("Helvetica", "B", 20)
    pdf.set_text_color(0, 200, 220)
    pdf.cell(0, 12, "MEETING INTELLIGENCE REPORT", ln=True, align="C")

    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(100, 110, 130)
    pdf.cell(0, 6, f"Role: {role.upper()}  |  File: {filename}  |  {datetime.now().strftime('%Y-%m-%d %H:%M')}", ln=True, align="C")
    pdf.ln(6)

    def section_header(title: str, r=0, g=220, b=200):
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(r, g, b)
        pdf.cell(0, 8, title, ln=True)
        pdf.set_draw_color(0, 180, 200)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(3)

    def body_text(text: str, color=(220, 230, 240)):
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(*color)
        pdf.multi_cell(0, 6, text)
        pdf.ln(2)

    # Summary
    section_header("📋 SUMMARY")
    body_text(result.summary)
    pdf.ln(2)

    # Key Themes
    if result.key_themes:
        section_header("🏷️ KEY THEMES", 180, 100, 255)
        for theme in result.key_themes:
            body_text(f"  • {theme}")
        pdf.ln(2)

    # Decisions
    section_header("🎯 KEY DECISIONS", 255, 80, 160)
    if result.decisions:
        for i, d in enumerate(result.decisions, 1):
            body_text(f"  {i}. {d}")
    else:
        body_text("  No explicit decisions identified.")
    pdf.ln(2)

    # Action Items
    section_header("⚡ ACTION ITEMS", 255, 200, 0)
    if result.action_items:
        for item in result.action_items:
            deadline_str = f" [by: {item.deadline}]" if item.deadline else ""
            priority_tag = f"[{item.priority.value.upper()}]"
            body_text(f"  {priority_tag} {item.task}\n       Owner: {item.owner}{deadline_str}")
    else:
        body_text("  No action items extracted.")
    pdf.ln(2)

    # Risks
    if result.risks:
        section_header("⚠️ RISKS & BLOCKERS", 255, 120, 0)
        for risk in result.risks:
            body_text(f"  ⚠ {risk}")

    return bytes(pdf.output())


def _build_pptx(result, role: str, filename: str) -> bytes:
    """Generate a PPTX deck from MeetingOutput using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise ImportError("python-pptx is required. Run: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK_BG   = RGBColor(10, 10, 18)
    CYAN      = RGBColor(0, 240, 255)
    MAGENTA   = RGBColor(255, 0, 153)
    MINT      = RGBColor(0, 255, 179)
    AMBER     = RGBColor(255, 200, 0)
    TEXT      = RGBColor(224, 230, 240)
    MUTED     = RGBColor(107, 115, 148)
    VIOLET    = RGBColor(139, 0, 255)

    blank_layout = prs.slide_layouts[6]  # completely blank

    def _add_slide(title_text: str, content_lines: list[str], title_color=CYAN,
                   content_color=TEXT, accent_color=CYAN):
        slide = prs.slides.add_slide(blank_layout)

        # Dark background rectangle
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()

        # Accent bar at top
        bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()

        # Title box
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.9))
        tf = tx_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.bold = True
        run.font.size = Pt(28)
        run.font.color.rgb = title_color

        # Divider
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(12.3), Emu(18000))
        line.fill.solid()
        line.fill.fore_color.rgb = accent_color
        line.line.fill.background()

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8))
        tf2 = content_box.text_frame
        tf2.word_wrap = True
        for i, line_text in enumerate(content_lines):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            run2 = p2.add_run()
            run2.text = line_text
            run2.font.size = Pt(15)
            run2.font.color.rgb = content_color
            p2.space_after = Pt(4)

        return slide

    # ── Slide 1: Title ──────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid(); bg1.fill.fore_color.rgb = DARK_BG; bg1.line.fill.background()

    bar1 = slide1.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.12))
    bar1.fill.solid(); bar1.fill.fore_color.rgb = CYAN; bar1.line.fill.background()

    bar2 = slide1.shapes.add_shape(1, 0, Inches(7.38), prs.slide_width, Inches(0.12))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = MAGENTA; bar2.line.fill.background()

    tb1 = slide1.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12), Inches(1.4))
    tf1 = tb1.text_frame; p1 = tf1.paragraphs[0]; r1 = p1.add_run()
    r1.text = "MEETING INTELLIGENCE REPORT"
    r1.font.bold = True; r1.font.size = Pt(38); r1.font.color.rgb = CYAN
    p1.alignment = PP_ALIGN.CENTER

    tb2 = slide1.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(12), Inches(0.6))
    tf2b = tb2.text_frame; p2b = tf2b.paragraphs[0]; r2b = p2b.add_run()
    r2b.text = f"Role: {role.upper()}  |  {filename}"
    r2b.font.size = Pt(16); r2b.font.color.rgb = MUTED
    p2b.alignment = PP_ALIGN.CENTER

    tb3 = slide1.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(12), Inches(0.5))
    tf3 = tb3.text_frame; p3 = tf3.paragraphs[0]; r3 = p3.add_run()
    r3.text = f"Generated {datetime.now().strftime('%B %d, %Y at %H:%M')}"
    r3.font.size = Pt(13); r3.font.color.rgb = MUTED
    p3.alignment = PP_ALIGN.CENTER

    # ── Slide 2: Summary ────────────────────────────────────────────
    _add_slide("📋  EXECUTIVE SUMMARY", [result.summary], title_color=CYAN, accent_color=CYAN)

    # ── Slide 3: Key Themes ─────────────────────────────────────────
    if result.key_themes:
        theme_lines = [f"  ◆  {t}" for t in result.key_themes]
        _add_slide("🏷️  KEY THEMES", theme_lines, title_color=VIOLET, content_color=TEXT, accent_color=VIOLET)

    # ── Slide 4: Key Decisions ──────────────────────────────────────
    decision_lines = [f"  {i+1}.  {d}" for i, d in enumerate(result.decisions)] if result.decisions else ["  No explicit decisions identified."]
    _add_slide("🎯  KEY DECISIONS", decision_lines, title_color=MAGENTA, content_color=TEXT, accent_color=MAGENTA)

    # ── Slide 5: Action Items ────────────────────────────────────────
    action_lines = []
    for item in result.action_items:
        deadline_str = f"  [by: {item.deadline}]" if item.deadline else ""
        action_lines.append(f"  [{item.priority.value.upper()}]  {item.task}")
        action_lines.append(f"           Owner: {item.owner}{deadline_str}")
        action_lines.append("")
    if not action_lines:
        action_lines = ["  No action items extracted."]
    _add_slide("⚡  ACTION ITEMS", action_lines, title_color=AMBER, content_color=TEXT, accent_color=AMBER)

    # ── Slide 6: Risks ──────────────────────────────────────────────
    if result.risks:
        risk_lines = [f"  ⚠  {r}" for r in result.risks]
        _add_slide("⚠️  RISKS & BLOCKERS", risk_lines, title_color=RGBColor(255, 120, 0), content_color=TEXT, accent_color=RGBColor(255, 120, 0))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ===================================================================== #
#  PAGE CONFIG
# ===================================================================== #
st.set_page_config(
    page_title="🎙️ Media Analyzer | Private Intelligence",
    page_icon="🎙️",
    layout="wide",
    initial_sidebar_state="expanded",
)


# ===================================================================== #
#  SESSION STATE INIT
# ===================================================================== #
if "analyzer_role" not in st.session_state:
    st.session_state.analyzer_role = "Engineering"
if "transcribed_text" not in st.session_state:
    st.session_state.transcribed_text = ""
if "last_result" not in st.session_state:
    st.session_state.last_result = None
if "last_filename" not in st.session_state:
    st.session_state.last_filename = ""


# ===================================================================== #
#  THREE.JS SCENE TEMPLATE
# ===================================================================== #
BG_TEMPLATE = (Path(__file__).resolve().parent.parent / "background.html").read_text(
    encoding="utf-8"
)


# ===================================================================== #
#  CYBERPUNK CSS  — improved buttons, export buttons, status cards
# ===================================================================== #
st.markdown(
    f"""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Orbitron:wght@400;700;900&family=JetBrains+Mono:wght@300;400;500&display=swap');

    :root {{
      --cyan:    #00f0ff;
      --magenta: #ff0099;
      --violet:  #8b00ff;
      --mint:    #00ffb3;
      --amber:   #ffc800;
      --orange:  #ff7800;
      --bg:      #0a0a12;
      --card:    rgba(10, 10, 18, 0.82);
      --border:  rgba(0, 240, 255, 0.15);
      --text:    #e0e6f0;
      --muted:   #6b7394;
    }}

    html, body, [data-testid="stAppViewContainer"],
    [data-testid="stApp"] {{
        background: var(--bg) !important;
        color: var(--text) !important;
        font-family: 'JetBrains Mono', monospace !important;
    }}

    header[data-testid="stHeader"] {{ background: transparent !important; }}
    [data-testid="stAppViewBlockContainer"] {{ background: transparent !important; }}
    #MainMenu, footer, [data-testid="stDecoration"] {{ display: none !important; }}
    [data-testid="stSidebarNav"] {{ display: none !important; }}

    .block-container {{
        padding-top: 1.5rem !important;
        max-width: 920px !important;
    }}

    /* ── 3D background iframe ────────────────────────────────────── */
    iframe[title="streamlit_components_v1.html"] {{
        position: fixed !important;
        top: 0 !important; left: 0 !important;
        width: 100vw !important; height: 100vh !important;
        z-index: 0 !important;
        pointer-events: none !important;
        border: none !important;
    }}
    [data-testid="stHtml"],
    [data-testid="element-container"]:has(iframe) {{
        position: fixed !important;
        top: 0 !important; left: 0 !important;
        width: 100vw !important; height: 100vh !important;
        z-index: 0 !important;
        overflow: visible !important;
        pointer-events: none !important;
    }}
    .block-container, [data-testid="stVerticalBlock"] {{
        position: relative !important;
        z-index: 1 !important;
    }}

    /* ── Typography ─────────────────────────────────────────────── */
    h1 {{
        font-family: 'Orbitron', sans-serif !important;
        font-weight: 900 !important;
        font-size: 1.8rem !important;
        background: linear-gradient(135deg, var(--cyan), var(--magenta), var(--violet));
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        background-clip: text;
        letter-spacing: 3px;
        text-transform: uppercase;
        padding-bottom: 0.2rem;
    }}
    .page-sub {{
        font-size: 0.75rem;
        color: var(--muted);
        letter-spacing: 4px;
        text-transform: uppercase;
        margin-bottom: 1.8rem;
    }}
    h3 {{
        font-family: 'Orbitron', sans-serif !important;
        font-size: 0.9rem !important; letter-spacing: 2px !important;
        color: var(--cyan) !important;
        border-bottom: 1px solid var(--border);
        padding-bottom: 0.5rem; margin-top: 1.5rem !important;
    }}

    /* ── Glass card ────────────────────────────────────────────── */
    .glass-card {{
        background: var(--card);
        backdrop-filter: blur(18px);
        -webkit-backdrop-filter: blur(18px);
        border: 1px solid var(--border);
        border-radius: 14px;
        padding: 1.8rem;
        margin-bottom: 1.5rem;
        box-shadow: 0 0 40px rgba(0,240,255,0.05), inset 0 0 60px rgba(0,240,255,0.02);
    }}

    /* ── Selectbox ──────────────────────────────────────────────── */
    .stSelectbox label, .stTextArea label, label {{
        font-family: 'Orbitron', sans-serif !important;
        font-size: 0.7rem !important;
        letter-spacing: 2px !important;
        text-transform: uppercase !important;
        color: var(--cyan) !important;
    }}
    [data-testid="stSelectbox"] > div > div,
    [data-testid="stSelectbox"] div[data-baseweb="select"] > div,
    [data-testid="stSelectbox"] [data-baseweb="select"],
    [data-testid="stSelectbox"] input {{
        background: rgba(0, 240, 255, 0.05) !important;
        background-color: rgba(0, 240, 255, 0.05) !important;
        border: 1px solid var(--border) !important;
        border-radius: 8px !important;
        color: var(--text) !important;
    }}
    [data-baseweb="popover"], [data-baseweb="popover"] ul,
    [data-baseweb="menu"], [role="listbox"] {{
        background: #0d0d1a !important;
        background-color: #0d0d1a !important;
        border: 1px solid var(--border) !important;
    }}
    [data-baseweb="menu"] li, [role="option"] {{ color: var(--text) !important; }}
    [data-baseweb="menu"] li:hover, [role="option"]:hover {{
        background: rgba(0, 240, 255, 0.1) !important;
    }}

    /* ── Textarea ───────────────────────────────────────────────── */
    .stTextArea textarea, textarea {{
        background: rgba(10, 10, 18, 0.9) !important;
        background-color: rgba(10, 10, 18, 0.9) !important;
        color: #e0e6f0 !important;
        caret-color: #00f0ff !important;
        border: 1px solid var(--border) !important;
        border-radius: 8px !important;
        font-family: 'JetBrains Mono', monospace !important;
        font-size: 0.82rem !important;
        line-height: 1.6 !important;
    }}
    .stTextArea textarea:focus, textarea:focus {{
        border-color: var(--cyan) !important;
        box-shadow: 0 0 15px rgba(0,240,255,0.2) !important;
        outline: none !important;
    }}

    /* ── PRIMARY ACTION BUTTONS (Transcribe / Analyze) ─────────── */
    .stButton > button {{
        width: 100%;
        font-family: 'Orbitron', sans-serif !important;
        font-weight: 700;
        font-size: 0.9rem;
        letter-spacing: 4px;
        text-transform: uppercase;
        color: #050510 !important;
        background: linear-gradient(135deg, #00f0ff 0%, #00ffb3 100%) !important;
        border: none !important;
        border-radius: 10px !important;
        padding: 0.85rem 2rem !important;
        cursor: pointer;
        transition: all 0.25s cubic-bezier(0.4, 0, 0.2, 1);
        box-shadow: 0 0 25px rgba(0,240,255,0.3), 0 4px 15px rgba(0,0,0,0.4);
        position: relative;
        overflow: hidden;
    }}
    .stButton > button::before {{
        content: '';
        position: absolute;
        top: 0; left: -100%;
        width: 100%; height: 100%;
        background: linear-gradient(90deg, transparent, rgba(255,255,255,0.15), transparent);
        transition: left 0.5s ease;
    }}
    .stButton > button:hover::before {{
        left: 100%;
    }}
    .stButton > button:hover {{
        box-shadow: 0 0 45px rgba(0,240,255,0.5), 0 0 80px rgba(0,255,179,0.2), 0 6px 20px rgba(0,0,0,0.5);
        transform: translateY(-2px);
        background: linear-gradient(135deg, #00f0ff 0%, #00ffc8 100%) !important;
    }}
    .stButton > button:active {{
        transform: translateY(0px);
        box-shadow: 0 0 20px rgba(0,240,255,0.4);
    }}

    /* BACK button override — smaller, ghost */
    [data-testid="stHorizontalBlock"]:first-of-type .stButton > button {{
        background: transparent !important;
        border: 1px solid var(--border) !important;
        color: var(--muted) !important;
        font-size: 0.7rem !important;
        padding: 0.5rem 1rem !important;
        letter-spacing: 2px !important;
        box-shadow: none !important;
        border-radius: 6px !important;
        width: auto !important;
    }}
    [data-testid="stHorizontalBlock"]:first-of-type .stButton > button:hover {{
        border-color: var(--cyan) !important;
        color: var(--cyan) !important;
        box-shadow: 0 0 10px rgba(0,240,255,0.15) !important;
        transform: none !important;
    }}

    /* ── EXPORT BUTTONS (PDF / PPTX) — secondary styling ───────── */
    .export-btn-row {{
        display: flex;
        gap: 1rem;
        margin-top: 1rem;
        margin-bottom: 0.5rem;
    }}

    /* Download buttons from st.download_button */
    [data-testid="stDownloadButton"] > button {{
        font-family: 'Orbitron', sans-serif !important;
        font-weight: 700 !important;
        font-size: 0.78rem !important;
        letter-spacing: 3px !important;
        text-transform: uppercase !important;
        color: var(--text) !important;
        background: rgba(10,10,18,0.9) !important;
        border: 1px solid var(--border) !important;
        border-radius: 8px !important;
        padding: 0.65rem 1.6rem !important;
        transition: all 0.25s ease !important;
        box-shadow: 0 0 15px rgba(0,240,255,0.08) !important;
        width: 100% !important;
    }}
    /* PDF button — amber accent */
    [data-testid="stDownloadButton"]:nth-of-type(1) > button {{
        border-color: rgba(255,200,0,0.3) !important;
    }}
    [data-testid="stDownloadButton"]:nth-of-type(1) > button:hover {{
        border-color: var(--amber) !important;
        color: var(--amber) !important;
        box-shadow: 0 0 20px rgba(255,200,0,0.25) !important;
        transform: translateY(-1px) !important;
    }}
    /* PPTX button — magenta accent */
    [data-testid="stDownloadButton"]:nth-of-type(2) > button {{
        border-color: rgba(255,0,153,0.3) !important;
    }}
    [data-testid="stDownloadButton"]:nth-of-type(2) > button:hover {{
        border-color: var(--magenta) !important;
        color: var(--magenta) !important;
        box-shadow: 0 0 20px rgba(255,0,153,0.25) !important;
        transform: translateY(-1px) !important;
    }}

    /* ── File uploader ──────────────────────────────────────────── */
    [data-testid="stFileUploader"] {{
        border: 1px dashed rgba(0,240,255,0.2) !important;
        border-radius: 10px !important;
        padding: 0.8rem !important;
        background: rgba(0,240,255,0.02) !important;
        transition: border-color 0.3s;
    }}
    [data-testid="stFileUploader"]:hover {{
        border-color: rgba(0,240,255,0.4) !important;
    }}
    [data-testid="stFileUploader"] section {{
        background: transparent !important;
    }}
    [data-testid="stFileUploader"] button {{
        color: var(--cyan) !important;
        border-color: var(--border) !important;
        background: rgba(0, 240, 255, 0.05) !important;
    }}

    /* ── Alert ──────────────────────────────────────────────────── */
    [data-testid="stAlert"] {{
        background: rgba(0, 240, 255, 0.06) !important;
        border: 1px solid rgba(0, 240, 255, 0.2) !important;
        border-left: 4px solid var(--cyan) !important;
        border-radius: 8px !important;
        color: var(--text) !important;
        font-size: 0.88rem;
    }}

    /* ── Expander ───────────────────────────────────────────────── */
    [data-testid="stExpander"] {{
        border: 1px solid var(--border) !important;
        border-radius: 10px !important;
        background: rgba(10, 10, 18, 0.5) !important;
    }}
    [data-testid="stExpander"] summary {{
        color: var(--cyan) !important;
        font-family: 'Orbitron', sans-serif !important;
        font-size: 0.72rem !important;
        letter-spacing: 1px !important;
    }}

    /* ── Decision item ──────────────────────────────────────────── */
    .decision-item {{
        padding: 0.6rem 1rem; margin-bottom: 0.4rem;
        background: rgba(139, 0, 255, 0.07);
        border-left: 3px solid var(--violet);
        border-radius: 0 8px 8px 0;
        font-size: 0.85rem;
        transition: background 0.2s;
    }}
    .decision-item:hover {{ background: rgba(139, 0, 255, 0.12); }}

    /* ── Risk item ──────────────────────────────────────────────── */
    .risk-item {{
        padding: 0.6rem 1rem; margin-bottom: 0.4rem;
        background: rgba(255, 120, 0, 0.07);
        border-left: 3px solid var(--orange);
        border-radius: 0 8px 8px 0;
        font-size: 0.85rem;
    }}

    /* ── Theme pill ─────────────────────────────────────────────── */
    .theme-pill {{
        display: inline-block;
        background: rgba(139,0,255,0.1);
        border: 1px solid rgba(139,0,255,0.3);
        border-radius: 50px;
        padding: 0.3rem 0.9rem;
        font-size: 0.72rem;
        letter-spacing: 1px;
        color: #c87fff;
        margin: 0.25rem 0.2rem;
    }}

    /* ── Table ──────────────────────────────────────────────────── */
    .stTable table, .stDataFrame table {{ width: 100%; border-collapse: separate; border-spacing: 0; }}
    .stTable thead th, .stDataFrame thead th {{
        background: rgba(0, 240, 255, 0.08) !important;
        color: var(--cyan) !important;
        font-family: 'Orbitron', sans-serif !important;
        font-size: 0.65rem !important; letter-spacing: 2px !important;
        text-transform: uppercase !important;
        border-bottom: 1px solid var(--border) !important; padding: 0.7rem 1rem !important;
    }}
    .stTable tbody td, .stDataFrame tbody td {{
        background: rgba(10, 10, 18, 0.6) !important;
        color: var(--text) !important; font-size: 0.82rem !important;
        border-bottom: 1px solid rgba(0,240,255,0.06) !important; padding: 0.6rem 1rem !important;
    }}
    .stTable tbody tr:hover td {{ background: rgba(0, 240, 255, 0.05) !important; }}

    /* ── Priority badges ────────────────────────────────────────── */
    td:last-child {{
        font-weight: 700;
    }}

    /* ── Status widget ──────────────────────────────────────────── */
    [data-testid="stStatusWidget"] {{
        border: 1px solid var(--border) !important;
        border-radius: 10px !important;
        background: rgba(10,10,18,0.7) !important;
    }}

    /* ── Media preview card ──────────────────────────────────────── */
    .preview-card {{
        background: rgba(0,240,255,0.03);
        border: 1px solid rgba(0,240,255,0.18);
        border-radius: 14px;
        padding: 1.4rem 1.6rem;
        margin: 0.8rem 0 1.2rem 0;
        position: relative;
        overflow: hidden;
    }}
    .preview-card::before {{
        content: '';
        position: absolute;
        top: 0; left: 0; right: 0;
        height: 2px;
        background: linear-gradient(90deg, var(--cyan), var(--magenta), var(--violet));
    }}
    .preview-badge {{
        display: inline-flex;
        align-items: center;
        gap: 0.4rem;
        background: rgba(0,240,255,0.08);
        border: 1px solid rgba(0,240,255,0.2);
        border-radius: 50px;
        padding: 0.25rem 0.8rem;
        font-size: 0.65rem;
        letter-spacing: 2px;
        text-transform: uppercase;
        color: var(--cyan);
        margin-bottom: 1rem;
    }}
    .preview-meta-row {{
        display: flex;
        gap: 1.5rem;
        flex-wrap: wrap;
        margin-top: 0.8rem;
    }}
    .preview-meta-item {{
        display: flex;
        flex-direction: column;
        gap: 0.15rem;
    }}
    .preview-meta-label {{
        font-size: 0.6rem;
        letter-spacing: 2px;
        text-transform: uppercase;
        color: var(--muted);
    }}
    .preview-meta-value {{
        font-size: 0.82rem;
        color: var(--text);
        font-weight: 500;
    }}
    .preview-thumb {{
        border-radius: 8px;
        border: 1px solid rgba(0,240,255,0.15);
        width: 100%;
        aspect-ratio: 16/9;
        object-fit: cover;
    }}
    .waveform-bar-container {{
        display: flex;
        align-items: flex-end;
        gap: 3px;
        height: 48px;
        margin: 0.5rem 0;
    }}
    /* Streamlit video/audio player theming */
    video {{
        border-radius: 10px !important;
        border: 1px solid rgba(0,240,255,0.2) !important;
        width: 100% !important;
        background: #050510 !important;
    }}
    audio {{
        width: 100% !important;
        border-radius: 8px !important;
        filter: invert(1) hue-rotate(180deg) saturate(1.5) brightness(0.9);
    }}

    /* ── Scrollbar ──────────────────────────────────────────────── */
    ::-webkit-scrollbar {{ width: 6px; }}
    ::-webkit-scrollbar-track {{ background: var(--bg); }}
    ::-webkit-scrollbar-thumb {{ background: rgba(0,240,255,0.2); border-radius: 3px; }}
    </style>
    """,
    unsafe_allow_html=True,
)


# ===================================================================== #
#  HEADER + BACK NAVIGATION
# ===================================================================== #
vault_record = render_vault_sidebar(current_role=st.session_state.analyzer_role)

bcol1, bcol2 = st.columns([1, 8])
with bcol1:
    if st.button("← BACK"):
        st.switch_page("app.py")

st.title("🎙️ Media Analyzer")
st.markdown(
    '<p class="page-sub">Upload audio/video · Gemma 4 extracts intelligence · Export to PDF or PPTX</p>',
    unsafe_allow_html=True,
)


# ===================================================================== #
#  ROLE SELECTOR + 3D BACKGROUND
# ===================================================================== #
st.markdown('<div class="glass-card">', unsafe_allow_html=True)

col1, col2 = st.columns([1, 2])

with col1:
    role = st.selectbox(
        "Stakeholder Lens",
        options=["Engineering", "Product", "Management"],
        index=["Engineering", "Product", "Management"].index(
            st.session_state.analyzer_role
        ),
        key="role_select",
    )
    if role != st.session_state.analyzer_role:
        st.session_state.analyzer_role = role
        st.rerun()

with col2:
    role_desc = {
        "Engineering": "🔧  Tech debt · Blockers · Architecture · Reliability",
        "Product":     "📦  Roadmap · Features · Customer impact · Scope",
        "Management":  "📊  Risks · Deadlines · Resource allocation · OKRs",
    }
    st.markdown(
        f'<div style="padding:0.5rem 0; color: var(--muted); font-size:0.75rem; letter-spacing:1px;">{role_desc[role]}</div>',
        unsafe_allow_html=True,
    )

# Inject 3D scene
BG_HTML = BG_TEMPLATE.replace("__ACTIVE_ROLE__", st.session_state.analyzer_role)
components.html(BG_HTML, height=600, scrolling=False)


# ===================================================================== #
#  INPUT: FILE UPLOAD
# ===================================================================== #
uploaded_file = st.file_uploader(
    "Upload a meeting recording",
    type=[ext.lstrip(".") for ext in ALL_EXTENSIONS],
    help="Video (.mp4 .mov .webm .mkv .avi) or Audio (.wav .mp3 .m4a .ogg .flac .aac .wma)",
)

if uploaded_file is not None:
    file_bytes = uploaded_file.read()
    uploaded_file.seek(0)  # reset so downstream reads work
    file_size_mb = len(file_bytes) / (1024 * 1024)
    ext = Path(uploaded_file.name).suffix.lower()
    is_video = ext in {".mp4", ".mov", ".webm", ".mkv", ".avi"}

    # ================================================================= #
    #  MEDIA PREVIEW CARD
    # ================================================================= #
    st.markdown(
        f"""
        <div class="preview-card">
          <div class="preview-badge">
            {'🎬 VIDEO FILE' if is_video else '🎵 AUDIO FILE'}
          </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    # ── Layout: thumbnail col | player + meta col ─────────────────────
    if is_video:
        thumb_col, player_col = st.columns([1, 2], gap="medium")
    else:
        player_col = st.container()
        thumb_col  = None

    # ── Video thumbnail (ffmpeg frame grab at 3 s) ───────────────────
    thumb_bytes: bytes | None = None
    if is_video:
        try:
            from core.audio import _get_ffmpeg_bin
            import subprocess
            ffmpeg_bin = _get_ffmpeg_bin()
            # Write video to temp file for ffmpeg
            tmp_fd_v, tmp_vid = tempfile.mkstemp(suffix=ext, prefix="thumb_src_")
            with os.fdopen(tmp_fd_v, "wb") as fv:
                fv.write(file_bytes)
            tmp_fd_t, tmp_thumb = tempfile.mkstemp(suffix=".jpg", prefix="thumb_")
            os.close(tmp_fd_t)
            res = subprocess.run(
                [
                    ffmpeg_bin, "-y",
                    "-ss", "00:00:03",       # seek to 3 s
                    "-i", tmp_vid,
                    "-frames:v", "1",
                    "-q:v", "2",
                    "-vf", "scale=480:-1",   # 480 px wide, auto height
                    tmp_thumb,
                ],
                capture_output=True, timeout=30,
            )
            if res.returncode == 0 and Path(tmp_thumb).exists():
                with open(tmp_thumb, "rb") as ft:
                    thumb_bytes = ft.read()
            os.remove(tmp_vid)
            os.remove(tmp_thumb)
        except Exception:
            thumb_bytes = None  # fall back gracefully

    with (thumb_col if is_video else st.container()):
        if is_video:
            if thumb_bytes:
                st.markdown(
                    '<p style="font-size:0.65rem;letter-spacing:2px;text-transform:uppercase;'
                    'color:var(--muted);margin-bottom:0.4rem;">📸 THUMBNAIL PREVIEW</p>',
                    unsafe_allow_html=True,
                )
                st.image(thumb_bytes, use_container_width=True)
            else:
                st.markdown(
                    '<div style="aspect-ratio:16/9;background:rgba(0,240,255,0.04);'
                    'border:1px dashed rgba(0,240,255,0.15);border-radius:8px;'
                    'display:flex;align-items:center;justify-content:center;'
                    'color:var(--muted);font-size:0.7rem;letter-spacing:1px;">'
                    '🎬 No thumbnail</div>',
                    unsafe_allow_html=True,
                )

    with player_col:
        st.markdown(
            '<p style="font-size:0.65rem;letter-spacing:2px;text-transform:uppercase;'
            'color:var(--muted);margin-bottom:0.4rem;">'
            + ('▶️ VIDEO PLAYER' if is_video else '🔊 AUDIO PLAYER') + '</p>',
            unsafe_allow_html=True,
        )
        if is_video:
            st.video(file_bytes)
        else:
            st.audio(file_bytes, format=f"audio/{ext.lstrip('.')}")

        # ── File metadata strip ──────────────────────────────────────
        # Detect duration using imageio-ffmpeg if available
        duration_str = "—"
        try:
            from core.audio import _get_ffmpeg_bin
            import subprocess, json as _json
            ffmpeg_bin = _get_ffmpeg_bin()
            ffprobe_path = str(Path(ffmpeg_bin).parent / "ffprobe.exe")
            if not Path(ffprobe_path).exists():
                ffprobe_path = ffmpeg_bin.replace("ffmpeg", "ffprobe")
            if Path(ffprobe_path).exists():
                tmp_fd_m, tmp_media = tempfile.mkstemp(suffix=ext, prefix="probe_")
                with os.fdopen(tmp_fd_m, "wb") as fm:
                    fm.write(file_bytes)
                probe = subprocess.run(
                    [ffprobe_path, "-v", "quiet", "-print_format", "json",
                     "-show_format", tmp_media],
                    capture_output=True, text=True, timeout=15,
                )
                os.remove(tmp_media)
                if probe.returncode == 0:
                    info = _json.loads(probe.stdout)
                    secs = float(info.get("format", {}).get("duration", 0))
                    if secs > 0:
                        m, s = divmod(int(secs), 60)
                        h, m = divmod(m, 60)
                        duration_str = f"{h:02d}:{m:02d}:{s:02d}" if h else f"{m:02d}:{s:02d}"
        except Exception:
            pass

        st.markdown(
            f"""
            <div class="preview-meta-row">
              <div class="preview-meta-item">
                <span class="preview-meta-label">Filename</span>
                <span class="preview-meta-value">{uploaded_file.name}</span>
              </div>
              <div class="preview-meta-item">
                <span class="preview-meta-label">Size</span>
                <span class="preview-meta-value">{file_size_mb:.2f} MB</span>
              </div>
              <div class="preview-meta-item">
                <span class="preview-meta-label">Duration</span>
                <span class="preview-meta-value" style="color:var(--cyan)">{duration_str}</span>
              </div>
              <div class="preview-meta-item">
                <span class="preview-meta-label">Pipeline</span>
                <span class="preview-meta-value" style="color:var(--mint)">ffmpeg → WAV → Whisper</span>
              </div>
            </div>
            """,
            unsafe_allow_html=True,
        )

    st.markdown("")  # spacer

    if st.button("🎙️  TRANSCRIBE", use_container_width=True, key="btn_transcribe"):
        # Write to a named temp file with correct suffix
        suffix = Path(uploaded_file.name).suffix
        tmp_fd, tmp_path = tempfile.mkstemp(suffix=suffix, prefix="upload_")
        try:
            with os.fdopen(tmp_fd, "wb") as f:
                f.write(file_bytes)  # use already-read bytes
        except Exception:
            os.close(tmp_fd)
            raise

        try:
            with st.status("🔮 Processing meeting recording …", expanded=True) as status:
                status.update(label="🎙️  Stage 1/2 — Converting & transcribing with Whisper …")
                st.write("⏳ Converting to WAV → running Whisper (local inference) …")

                processor = get_audio_processor()
                transcript_text = processor.transcribe_file(tmp_path)

                st.write(f"✅ Transcribed {len(transcript_text):,} characters")
                status.update(label="✅ Transcription complete!", state="complete")

            st.session_state.transcribed_text = transcript_text
            st.session_state.last_filename = uploaded_file.name

        except Exception as exc:
            st.error(f"❌ Transcription failed: {exc}")
        finally:
            try:
                os.remove(tmp_path)
            except OSError:
                pass

# Show transcript preview
if st.session_state.transcribed_text:
    with st.expander("📜 Show Full Transcript", expanded=False):
        st.text(st.session_state.transcribed_text)


# ===================================================================== #
#  ANALYZE BUTTON
# ===================================================================== #
transcript = st.session_state.transcribed_text

process = st.button("⚡  ANALYZE WITH GEMMA 4", use_container_width=True, key="btn_analyze")
st.markdown("</div>", unsafe_allow_html=True)


# ===================================================================== #
#  ARCHIVED RECORD DISPLAY (from vault sidebar)
# ===================================================================== #
if vault_record is not None:
    st.markdown(
        '<div style="font-size:0.7rem; color:#00ffd4; letter-spacing:2px; '
        'text-transform:uppercase; margin-bottom:0.8rem;">'
        '🗄️ Viewing Archived Analysis</div>',
        unsafe_allow_html=True,
    )

    st.markdown("### 📋 Summary")
    st.info(vault_record.get("summary", ""))

    st.markdown("### 🎯 Key Decisions")
    decisions = vault_record.get("decisions", [])
    if decisions:
        for d in decisions:
            st.markdown(f'<div class="decision-item">{d}</div>', unsafe_allow_html=True)
    else:
        st.caption("No explicit decisions were identified.")

    st.markdown("### ⚡ Action Items")
    items = vault_record.get("action_items", [])
    if items:
        df = pd.DataFrame(items)
        col_map = {"task": "Task", "owner": "Owner", "priority": "Priority"}
        df.rename(columns={k: v for k, v in col_map.items() if k in df.columns}, inplace=True)
        if "Priority" in df.columns:
            df["Priority"] = df["Priority"].str.upper()
        st.table(df)
    else:
        st.caption("No action items were extracted.")

    with st.expander("📜 Transcript Snippet", expanded=False):
        st.text(vault_record.get("transcript_snippet", "(not available)"))

    st.stop()


# ===================================================================== #
#  PROCESSING & RESULTS
# ===================================================================== #
if process:
    if not transcript or not transcript.strip():
        st.warning("⚠️  Upload and transcribe a media file first.")
    else:
        transcript = transcript.strip()
        analyzer = get_analyzer()

        with st.status("🧠 Extracting intelligence with Gemma 4 …", expanded=True) as status:
            status.update(label="🧠  Stage 2/2 — Gemma 4 extracting insights …")
            st.write(f"🔍 Role: **{role}** · Sending {len(transcript):,} chars to local LLM …")

            try:
                result = analyzer.analyze(transcript, role)
            except RuntimeError as exc:
                st.error(f"❌  Inference failed: {exc}")
                st.stop()

            status.update(label="✅ Gemma 4 analysis complete!", state="complete")

        # ── Save to vault ────────────────────────────────────────────
        action_dicts = [
            {
                "task": item.task,
                "owner": item.owner,
                "priority": item.priority.value,
                "deadline": item.deadline,
            }
            for item in result.action_items
        ]
        save_analysis(
            department=role,
            summary=result.summary,
            decisions=result.decisions,
            action_items=action_dicts,
            transcript_snippet=transcript,
            source="media",
        )
        st.toast("💾 Saved to Meeting Vault", icon="✅")
        st.session_state.last_result = result

# ── Display results (from last_result in session state) ──────────────
result = st.session_state.last_result
if result is not None:
    filename = st.session_state.last_filename or "meeting"

    # ── Summary ──────────────────────────────────────────────────
    st.markdown("### 📋 Summary")
    st.info(result.summary)

    # ── Key Themes ───────────────────────────────────────────────
    if result.key_themes:
        st.markdown("### 🏷️ Key Themes")
        theme_html = "".join(f'<span class="theme-pill">◆ {t}</span>' for t in result.key_themes)
        st.markdown(f'<div style="margin-bottom:0.5rem;">{theme_html}</div>', unsafe_allow_html=True)

    # ── Decisions ────────────────────────────────────────────────
    st.markdown("### 🎯 Key Decisions")
    if result.decisions:
        for d in result.decisions:
            st.markdown(f'<div class="decision-item">{d}</div>', unsafe_allow_html=True)
    else:
        st.caption("No explicit decisions were identified.")

    # ── Action Items ─────────────────────────────────────────────
    st.markdown("### ⚡ Action Items")
    if result.action_items:
        rows = []
        for item in result.action_items:
            rows.append({
                "Task": item.task,
                "Owner": item.owner,
                "Priority": item.priority.value.upper(),
                "Deadline": item.deadline or "—",
            })
        df = pd.DataFrame(rows)
        st.table(df)
    else:
        st.caption("No action items were extracted.")

    # ── Risks ─────────────────────────────────────────────────────
    if result.risks:
        st.markdown("### ⚠️ Risks & Blockers")
        for r in result.risks:
            st.markdown(f'<div class="risk-item">⚠ {r}</div>', unsafe_allow_html=True)

    # ── Export Section ────────────────────────────────────────────
    st.markdown("---")
    st.markdown(
        '<p style="font-family:Orbitron,sans-serif; font-size:0.8rem; '
        'letter-spacing:3px; color:var(--cyan); text-transform:uppercase; margin-bottom:0.8rem;">'
        '📤 Export Summary</p>',
        unsafe_allow_html=True,
    )

    export_col1, export_col2 = st.columns(2)

    with export_col1:
        try:
            pdf_bytes = _build_pdf(result, role, filename)
            st.download_button(
                label="📄  DOWNLOAD PDF",
                data=pdf_bytes,
                file_name=f"meeting_summary_{role.lower()}_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf",
                mime="application/pdf",
                use_container_width=True,
                key="dl_pdf",
            )
        except ImportError as e:
            st.caption(f"PDF export requires: `pip install fpdf2`")
        except Exception as e:
            st.caption(f"PDF error: {e}")

    with export_col2:
        try:
            pptx_bytes = _build_pptx(result, role, filename)
            st.download_button(
                label="📊  DOWNLOAD PPTX",
                data=pptx_bytes,
                file_name=f"meeting_deck_{role.lower()}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                key="dl_pptx",
            )
        except ImportError as e:
            st.caption(f"PPTX export requires: `pip install python-pptx`")
        except Exception as e:
            st.caption(f"PPTX error: {e}")
